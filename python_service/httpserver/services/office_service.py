"""
Office Document Parsing Service.

Provides parsing capabilities for Office documents:
- XLSX/XLS (Excel)
- PPTX/PPT (PowerPoint)

Returns content as Markdown text for forensic analysis.
"""

import asyncio
import logging
import subprocess
from pathlib import Path
from typing import Optional

logger = logging.getLogger(__name__)


class OfficeService:
    """Service for parsing Office documents."""

    def __init__(self):
        """Initialize the Office parsing service."""
        pass

    async def parse_file(self, file_path: str) -> str:
        """
        Parse an Office file and return its content as Markdown.

        Args:
            file_path: Absolute path to the Office file.

        Returns:
            Extracted text content in Markdown format.

        Raises:
            ValueError: If file type is not supported.
            FileNotFoundError: If file does not exist.
        """
        path = Path(file_path)

        if not path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")

        suffix = path.suffix.lower()

        # The concrete parsers (openpyxl / python-pptx / subprocess) are
        # synchronous and can take seconds on large files; offload them to a
        # worker thread so the async event loop is not blocked.
        if suffix == ".xlsx":
            return await asyncio.to_thread(self._parse_xlsx, file_path)
        elif suffix == ".xls":
            return await asyncio.to_thread(self._parse_xls, file_path)
        elif suffix == ".pptx":
            return await asyncio.to_thread(self._parse_pptx, file_path)
        elif suffix == ".ppt":
            return await asyncio.to_thread(self._parse_ppt, file_path)
        else:
            raise ValueError(f"Unsupported file type: {suffix}")

    def _parse_xlsx(self, file_path: str) -> str:
        """Parse XLSX file using openpyxl."""
        try:
            from openpyxl import load_workbook

            wb = load_workbook(file_path, read_only=True, data_only=True)
            result = []

            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                result.append(f"## Sheet: {sheet_name}\n")

                # Build table rows
                rows = []
                for row in sheet.iter_rows():
                    cells = []
                    for cell in row:
                        value = cell.value if cell.value is not None else ""
                        cells.append(str(value).replace("|", "\\|"))
                    if any(cells):  # Skip empty rows
                        rows.append(cells)

                if rows:
                    # Create markdown table
                    if rows:
                        # Header row
                        result.append("| " + " | ".join(rows[0]) + " |")
                        result.append("|" + "|".join(["---"] * len(rows[0])) + "|")
                        # Data rows
                        for row in rows[1:]:
                            # Pad row to match header length
                            while len(row) < len(rows[0]):
                                row.append("")
                            result.append("| " + " | ".join(row[:len(rows[0])]) + " |")
                    result.append("")

            wb.close()
            return "\n".join(result)

        except Exception as e:
            logger.error(f"Error parsing XLSX {file_path}: {e}")
            return f"Error parsing XLSX file: {e}"

    def _parse_xls(self, file_path: str) -> str:
        """Parse XLS file using xlrd."""
        try:
            import xlrd

            wb = xlrd.open_workbook(file_path)
            result = []

            for sheet_idx in range(wb.nsheets):
                sheet = wb.sheet_by_index(sheet_idx)
                result.append(f"## Sheet: {sheet.name}\n")

                rows = []
                for row_idx in range(sheet.nrows):
                    cells = []
                    for col_idx in range(sheet.ncols):
                        value = sheet.cell_value(row_idx, col_idx)
                        cells.append(str(value).replace("|", "\\|"))
                    if any(cells):
                        rows.append(cells)

                if rows:
                    # Create markdown table
                    result.append("| " + " | ".join(rows[0]) + " |")
                    result.append("|" + "|".join(["---"] * len(rows[0])) + "|")
                    for row in rows[1:]:
                        while len(row) < len(rows[0]):
                            row.append("")
                        result.append("| " + " | ".join(row[:len(rows[0])]) + " |")
                    result.append("")

            return "\n".join(result)

        except Exception as e:
            logger.error(f"Error parsing XLS {file_path}: {e}")
            return f"Error parsing XLS file: {e}"

    def _parse_pptx(self, file_path: str) -> str:
        """Parse PPTX file using python-pptx."""
        try:
            from pptx import Presentation

            prs = Presentation(file_path)
            result = []

            for slide_num, slide in enumerate(prs.slides, 1):
                result.append(f"## Slide {slide_num}\n")

                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        # Check if it's a title
                        if shape.is_placeholder and hasattr(shape, "placeholder_format"):
                            ph_type = shape.placeholder_format.type
                            # Title placeholder types: 1 = CENTER_TITLE, 3 = TITLE
                            if ph_type in [1, 3]:
                                result.append(f"### {shape.text.strip()}\n")
                            else:
                                result.append(f"{shape.text.strip()}\n")
                        else:
                            result.append(f"{shape.text.strip()}\n")

                    # Handle tables in slides
                    if shape.has_table:
                        table = shape.table
                        table_rows = []
                        for row in table.rows:
                            cells = []
                            for cell in row.cells:
                                text = cell.text.replace("|", "\\|").replace("\n", " ")
                                cells.append(text)
                            table_rows.append(cells)

                        if table_rows:
                            result.append("| " + " | ".join(table_rows[0]) + " |")
                            result.append("|" + "|".join(["---"] * len(table_rows[0])) + "|")
                            for row in table_rows[1:]:
                                while len(row) < len(table_rows[0]):
                                    row.append("")
                                result.append("| " + " | ".join(row[:len(table_rows[0])]) + " |")
                            result.append("")

                result.append("")

            return "\n".join(result)

        except Exception as e:
            logger.error(f"Error parsing PPTX {file_path}: {e}")
            return f"Error parsing PPTX file: {e}"

    def _parse_ppt(self, file_path: str) -> str:
        """Parse PPT file using catppt (from catdoc package)."""
        try:
            result = subprocess.run(
                ["catppt", file_path],
                capture_output=True,
                text=True,
                timeout=60
            )

            if result.returncode != 0:
                logger.warning(f"catppt error: {result.stderr}")
                return f"Error parsing PPT: {result.stderr}"

            # Format output as markdown
            lines = result.stdout.strip().split("\n")
            md_lines = []
            current_slide = 0

            for line in lines:
                if line.startswith("Slide"):
                    current_slide += 1
                    md_lines.append(f"\n## Slide {current_slide}\n")
                elif line.strip():
                    md_lines.append(line.strip())

            return "\n".join(md_lines) if md_lines else result.stdout

        except FileNotFoundError:
            logger.error("catppt not found. Install catdoc package.")
            return "Error: catppt not found. Please install catdoc package."
        except subprocess.TimeoutExpired:
            logger.error(f"Timeout parsing PPT {file_path}")
            return "Error: Timeout parsing PPT file."
        except Exception as e:
            logger.error(f"Error parsing PPT {file_path}: {e}")
            return f"Error parsing PPT file: {e}"


# Global service instance
_office_service: Optional[OfficeService] = None


def get_office_service() -> OfficeService:
    """Get the global Office service instance."""
    global _office_service
    if _office_service is None:
        _office_service = OfficeService()
    return _office_service
